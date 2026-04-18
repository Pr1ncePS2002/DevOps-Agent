"""Generate AI DevOps Commander project presentation (20-24 slides)."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)
ACCENT  = RGBColor(0x38, 0xBD, 0xF8)
PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
GREEN   = RGBColor(0x4A, 0xDE, 0x80)
ORANGE  = RGBColor(0xFB, 0xBF, 0x24)
RED     = RGBColor(0xF8, 0x71, 0x71)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x94, 0xA3, 0xB8)
LIGHT   = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────
def bg(slide, color=BG_DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, border=None, radius=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    s.adjustments[0] = radius
    return s


def line(slide, l, t, w, color=ACCENT, thick=Pt(3)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, thick)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h).text_frame


def txt(tf, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        after=Pt(6), before=Pt(0), font="Segoe UI"):
    if tf.paragraphs[0].text == "" and not hasattr(tf, '_used'):
        p = tf.paragraphs[0]; tf._used = True
    else:
        p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font; p.alignment = align
    p.space_after = after; p.space_before = before
    return p


def header(slide, num, title, subtitle=None):
    bg(slide)
    badge = rect(slide, Inches(0.6), Inches(0.4), Inches(0.7), Inches(0.55), ACCENT)
    btf = badge.text_frame; btf.word_wrap = False
    bp = btf.paragraphs[0]; bp.text = num; bp.font.size = Pt(20)
    bp.font.color.rgb = BG_DARK; bp.font.bold = True; bp.font.name = "Segoe UI"
    bp.alignment = PP_ALIGN.CENTER; btf.vertical_anchor = MSO_ANCHOR.MIDDLE

    tf = tb(slide, Inches(1.5), Inches(0.35), Inches(10), Inches(0.6))
    txt(tf, title, sz=30, bold=True)
    line(slide, Inches(0.6), Inches(1.05), Inches(11.5))
    if subtitle:
        tf2 = tb(slide, Inches(0.6), Inches(1.15), Inches(11), Inches(0.45))
        txt(tf2, subtitle, sz=15, color=GRAY)


def bullets(slide, items, l, t, w, h, sz=16, col=LIGHT, bcol=ACCENT, sp=Pt(8)):
    tf = tb(slide, l, t, w, h); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = sp; p.space_before = Pt(2)
        r1 = p.add_run(); r1.text = "\u25B8  "; r1.font.size = Pt(sz)
        r1.font.color.rgb = bcol; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = item; r2.font.size = Pt(sz)
        r2.font.color.rgb = col; r2.font.name = "Segoe UI"


def card(slide, l, t, w, h, title, items, ac=ACCENT, tsz=16, isz=14):
    rect(slide, l, t, w, h, BG_CARD, border=ac)
    tf = tb(slide, l + Inches(0.2), t + Inches(0.15), w - Inches(0.4), Inches(0.4))
    txt(tf, title, sz=tsz, bold=True, color=ac)
    bullets(slide, items, l + Inches(0.2), t + Inches(0.55),
            w - Inches(0.4), h - Inches(0.7), sz=isz, col=LIGHT, bcol=GRAY, sp=Pt(5))


def card_text(slide, l, t, w, h, title, desc, ac=ACCENT):
    rect(slide, l, t, w, h, BG_CARD, border=ac)
    tf = tb(slide, l + Inches(0.12), t + Inches(0.1), w - Inches(0.24), Inches(0.35))
    txt(tf, title, sz=14, bold=True, color=ac, align=PP_ALIGN.CENTER)
    tf2 = tb(slide, l + Inches(0.12), t + Inches(0.5), w - Inches(0.24), h - Inches(0.6))
    tf2.word_wrap = True
    txt(tf2, desc, sz=12, color=LIGHT, align=PP_ALIGN.CENTER)


def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT)
rect(s, Inches(1.5), Inches(1.2), Inches(10.3), Inches(4.8), BG_CARD,
     border=RGBColor(0x2A, 0x3A, 0x5C))

tf = tb(s, Inches(2), Inches(1.5), Inches(9.3), Inches(1.4))
tf.word_wrap = True
txt(tf, "AI DEVOPS COMMANDER", sz=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(tf, "Natural Language Driven Deployment Automation", sz=22, color=LIGHT, align=PP_ALIGN.CENTER, after=Pt(2))
txt(tf, "with Agentic AI", sz=22, color=LIGHT, align=PP_ALIGN.CENTER)

line(s, Inches(4.5), Inches(3.6), Inches(4.3), thick=Pt(2))

tf2 = tb(s, Inches(2), Inches(3.85), Inches(9.3), Inches(1.8))
tf2.word_wrap = True
txt(tf2, "Submitted by", sz=14, color=GRAY, align=PP_ALIGN.CENTER)
txt(tf2, "Prince Parmar Singh  (24100443)", sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, after=Pt(14))
txt(tf2, "Under the guidance of", sz=14, color=GRAY, align=PP_ALIGN.CENTER)
txt(tf2, "Dr. Deepika MP,  Assistant Professor", sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tf3 = tb(s, Inches(2), Inches(6.2), Inches(9.3), Inches(0.7))
tf3.word_wrap = True
txt(tf3, "Department of Computer Applications", sz=13, color=GRAY, align=PP_ALIGN.CENTER, after=Pt(2))
txt(tf3, "Cochin University of Science and Technology (CUSAT)  |  April 2026", sz=13, color=GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ACCENT)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Introduction
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "01", "Introduction")

card(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.4),
     "The Deployment Bottleneck", [
         "Modern CI/CD demands deep CLI, YAML, and cloud-provider expertise",
         "Deployment errors cost engineering teams hours of debugging",
         "Non-DevOps developers are locked out of deployment workflows",
         "Monitoring, rollback, and logs are fragmented across tools",
     ], ac=ORANGE)

card(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.4),
     "The LLM Opportunity", [
         "Large Language Models can reason about structured plans from text",
         "Conversational AI can bridge intent and infrastructure",
         "Agentic pipelines enable safe, multi-step automation",
         "Natural language lowers the barrier for all team members",
     ], ac=GREEN)

card(s, Inches(0.6), Inches(4.3), Inches(12), Inches(2.6),
     "What is AI DevOps Commander?", [
         "A full-stack agentic system: users deploy software using plain English commands",
         "Next.js 14 dashboard  +  FastAPI backend  +  Redis job queue  +  Docker SDK execution",
         "Multi-strategy interpreter: LLM primary (Google Gemini) with deterministic regex fallback",
         "LLM-powered risk assessment, mandatory approval gates, automatic watchdog rollback",
         "Layered safety: dry-run default, execution gates, path whitelists, encrypted secrets",
     ], ac=ACCENT, isz=15)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "02", "Problem Statement")

problems = [
    ("Technical Barriers", "Deployment requires Docker CLI, CI/CD YAML, cloud APIs, and shell scripting \u2014 creating an access gap for non-infrastructure specialists.", RED),
    ("Cognitive Overload", "Managing multiple environments, tracking versions, and coordinating rollbacks simultaneously places unsustainable cognitive demands.", ORANGE),
    ("No Contextual Intelligence", "Existing tools execute commands blindly \u2014 deploying to production on a Friday without tests is perfectly allowed.", ORANGE),
    ("Fragmented Observability", "Status, health, rollback triggers, and logs are scattered across multiple dashboards and CLIs.", ACCENT),
    ("No Conversational Interface", "No current deployment tool supports natural language interaction for infrastructure operations.", PURPLE),
]

for i, (title, desc, color) in enumerate(problems):
    col = i % 3; row = i // 3
    left = Inches(0.6) + col * Inches(4.15)
    top = Inches(1.6) + row * Inches(2.7)
    rect(s, left, top, Inches(3.9), Inches(2.4), BG_CARD, border=color)
    tf = tb(s, left + Inches(0.2), top + Inches(0.15), Inches(3.5), Inches(0.4))
    txt(tf, title, sz=16, bold=True, color=color)
    tf2 = tb(s, left + Inches(0.2), top + Inches(0.6), Inches(3.5), Inches(1.6))
    tf2.word_wrap = True
    txt(tf2, desc, sz=13, color=LIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Objectives
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "03", "Objectives \u2014 Four Key Pillars")

pillars = [
    ("Natural Language\nInterface", ACCENT, [
        "Accept deployment commands in plain English",
        "Multi-turn chat with intent detection",
        "No CLI syntax knowledge needed",
    ]),
    ("Intelligent\nRisk Assessment", ORANGE, [
        "LLM-powered plan advisor (risk 0\u2013100)",
        "Env, test coverage, version analysis",
        "Deterministic fallback for reliability",
    ]),
    ("Safe Automated\nExecution", GREEN, [
        "6-stage Docker pipeline via Python SDK",
        "Mandatory approval gate before execution",
        "Dry-run + execution gate + path whitelist",
    ]),
    ("Self-Healing\nInfrastructure", PURPLE, [
        "60s watchdog monitors container health",
        "Auto-rollback to last-known-good image",
        "Encrypted env vars, structured logging",
    ]),
]

for i, (title, color, items) in enumerate(pillars):
    left = Inches(0.5) + i * Inches(3.15)
    rect(s, left, Inches(1.6), Inches(2.95), Inches(5.2), BG_CARD, border=color)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.1), Inches(1.8), Inches(0.65), Inches(0.65))
    circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
    ctf = circle.text_frame; cp = ctf.paragraphs[0]
    cp.text = str(i + 1); cp.font.size = Pt(22); cp.font.bold = True
    cp.font.color.rgb = BG_DARK; cp.alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = tb(s, left + Inches(0.15), Inches(2.6), Inches(2.65), Inches(0.8))
    tf.word_wrap = True
    txt(tf, title, sz=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    bullets(s, items, left + Inches(0.15), Inches(3.5), Inches(2.65), Inches(3),
            sz=13, col=LIGHT, bcol=GRAY, sp=Pt(6))


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Proposed System: Request Flow
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "04", "Proposed System \u2014 End-to-End Request Flow")

steps = [
    ("\u2460 User Input", "Plain English\ncommand via chat\nor console", ACCENT),
    ("\u2461 Interpreter", "LLM parse with\nregex fallback\n(confidence gate)", PURPLE),
    ("\u2462 Risk Advisor", "RAG-style risk\nscoring (0\u2013100)\n+ warnings", ORANGE),
    ("\u2463 Approval", "Mandatory human\nreview & explicit\napproval", GREEN),
    ("\u2464 Orchestrator", "6-stage Docker\npipeline execution\n(build\u2192run\u2192check)", ACCENT),
    ("\u2465 Watchdog", "60s health poll\nauto-rollback on\ncontainer failure", RED),
]

bw, bh = Inches(1.75), Inches(1.8)
gap = Inches(0.22)
start_x = Inches(0.35)
y = Inches(2.0)

for i, (lbl, desc, col) in enumerate(steps):
    x = start_x + i * (bw + gap)
    rect(s, x, y, bw, bh, BG_CARD, border=col)
    tf = tb(s, x + Inches(0.05), y + Inches(0.1), bw - Inches(0.1), Inches(0.35))
    txt(tf, lbl, sz=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    tf2 = tb(s, x + Inches(0.05), y + Inches(0.45), bw - Inches(0.1), Inches(1.2))
    tf2.word_wrap = True
    txt(tf2, desc, sz=11, color=LIGHT, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        ax = x + bw
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, y + Inches(0.65), gap, Inches(0.35))
        arr.fill.solid(); arr.fill.fore_color.rgb = GRAY; arr.line.fill.background()

# Bottom: architecture layers
layers = [
    ("Frontend Layer", "Next.js 14  \u2022  React 18  \u2022  Tailwind CSS  \u2022  TypeScript  \u2022  Port 3000", ACCENT),
    ("Backend Layer", "FastAPI  \u2022  SQLModel  \u2022  Pydantic v2  \u2022  structlog  \u2022  Port 3001", GREEN),
    ("Infrastructure", "Docker SDK  \u2022  Redis + RQ  \u2022  SQLite  \u2022  Encrypted Secrets  \u2022  Watchdog", ORANGE),
    ("LLM Layer", "Google Gemini (primary)  \u2022  OpenAI  \u2022  Ollama  \u2022  Provider Abstraction", PURPLE),
]

for i, (lbl, desc, col) in enumerate(layers):
    top = Inches(4.3) + i * Inches(0.78)
    rect(s, Inches(0.6), top, Inches(12), Inches(0.65), BG_CARD, border=col)
    tf = tb(s, Inches(0.8), top + Inches(0.08), Inches(2.2), Inches(0.5))
    txt(tf, lbl, sz=13, bold=True, color=col)
    tf2 = tb(s, Inches(3.2), top + Inches(0.08), Inches(9), Inches(0.5))
    txt(tf2, desc, sz=13, color=LIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Feasibility Analysis
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "05", "Feasibility Analysis")

feas = [
    ("Technical Feasibility", GREEN, [
        "Built on mature, production-grade open-source frameworks",
        "FastAPI, Next.js, Docker SDK, Redis, SQLModel \u2014 all well-documented",
        "Google Gemini provides reliable LLM API with structured output",
        "Requires Python 3.11+, Node.js 18+, Docker \u2014 widely available",
    ]),
    ("Operational Feasibility", ACCENT, [
        "NL interface lowers barrier for non-DevOps users",
        "Risk advisor provides intelligent guardrails",
        "Orchestrator delivers reliable container management",
        "Watchdog provides automated recovery without human intervention",
        "Validated through unit, integration, and E2E testing",
    ]),
    ("Economic Feasibility", ORANGE, [
        "Entirely open-source \u2014 zero proprietary license costs",
        "LLM API costs minimal for development and demo use",
        "Only infrastructure: a machine with Docker installed",
        "Redis can run as a Docker container itself",
        "Full Docker Compose support for one-command deployment",
    ]),
]

for i, (title, color, items) in enumerate(feas):
    left = Inches(0.5) + i * Inches(4.15)
    card(s, left, Inches(1.6), Inches(3.95), Inches(5.2), title, items,
         ac=color, tsz=17, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Technologies Used (1/2: Core)
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "06", "Technologies Used \u2014 Core Stack")

# Left: Frontend
card(s, Inches(0.6), Inches(1.5), Inches(3.8), Inches(5.3),
     "Frontend", [
         "Next.js 14 \u2014 React framework",
         "React 18 \u2014 Component library",
         "TypeScript 5.5 \u2014 Type safety",
         "Tailwind CSS 3.4 \u2014 Utility styling",
         "Vitest 2.0 \u2014 Unit testing",
         "lucide-react \u2014 Icon library",
     ], ac=ACCENT, isz=15)

# Middle: Backend
card(s, Inches(4.75), Inches(1.5), Inches(3.8), Inches(5.3),
     "Backend", [
         "Python 3.12 \u2014 Runtime",
         "FastAPI 0.115 \u2014 API framework",
         "SQLModel 0.0.22 \u2014 ORM (SQLAlchemy)",
         "Pydantic v2 \u2014 Data validation",
         "RQ + Redis \u2014 Job queue",
         "structlog \u2014 Structured logging",
         "pytest + anyio \u2014 Testing",
     ], ac=GREEN, isz=15)

# Right: Infrastructure & AI
card(s, Inches(8.9), Inches(1.5), Inches(3.8), Inches(5.3),
     "Infrastructure & AI", [
         "Docker Python SDK \u2014 Container ops",
         "Redis 7 \u2014 Job queue backend",
         "SQLite \u2014 Dev database",
         "cryptography \u2014 Secret encryption",
         "Google Gemini \u2014 Primary LLM",
         "OpenAI / Ollama \u2014 Alternate LLMs",
         "httpx \u2014 Async HTTP client",
     ], ac=PURPLE, isz=15)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — System Design: Architecture
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "07", "System Design \u2014 Service Architecture")

# Backend services table
services = [
    ("command_interpreter.py", "Parses NL commands; LLM primary, regex fallback", ACCENT),
    ("orchestrator.py", "Executes approved deployment plans (6-stage pipeline)", GREEN),
    ("rag_advisor.py", "LLM-powered risk assessment with deterministic fallback", ORANGE),
    ("watchdog.py", "Polls container health; triggers auto-rollback", RED),
    ("project_analysis.py", "Detects tech stack; auto-generates Dockerfiles", PURPLE),
    ("conversation_engine.py", "Multi-turn chat with deployment intent detection", ACCENT),
    ("llm/factory.py", "Provider-agnostic LLM client factory", GRAY),
]

# Header row
rect(s, Inches(0.6), Inches(1.5), Inches(4), Inches(0.45), ACCENT)
tf = tb(s, Inches(0.8), Inches(1.5), Inches(3.6), Inches(0.45))
txt(tf, "Service Module", sz=14, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
rect(s, Inches(4.6), Inches(1.5), Inches(7.6), Inches(0.45), ACCENT)
tf = tb(s, Inches(4.8), Inches(1.5), Inches(7.2), Inches(0.45))
txt(tf, "Responsibility", sz=14, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

for i, (mod, role, col) in enumerate(services):
    top = Inches(2.05) + i * Inches(0.55)
    rect(s, Inches(0.6), top, Inches(4), Inches(0.48), BG_CARD, border=col)
    tf = tb(s, Inches(0.8), top + Inches(0.05), Inches(3.6), Inches(0.4))
    txt(tf, mod, sz=13, bold=True, color=col)
    rect(s, Inches(4.6), top, Inches(7.6), Inches(0.48), BG_CARD)
    tf2 = tb(s, Inches(4.8), top + Inches(0.05), Inches(7.2), Inches(0.4))
    txt(tf2, role, sz=13, color=LIGHT)

# API Routes
rect(s, Inches(0.6), Inches(6.0), Inches(12), Inches(1.0), BG_CARD, border=ACCENT)
tf = tb(s, Inches(0.8), Inches(6.05), Inches(11.5), Inches(0.35))
txt(tf, "API Routes", sz=14, bold=True, color=ACCENT)
tf2 = tb(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5))
txt(tf2, "projects  \u2022  commands  \u2022  executions  \u2022  providers  \u2022  deploy_status  \u2022  demo  \u2022  chat",
    sz=14, color=LIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — System Design: Use Case Diagram
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "07", "System Design \u2014 Use Case Diagram")

# Actors
actors_data = [
    ("User", "Developer / Operator\nvia web dashboard", ACCENT),
    ("LLM Provider", "Gemini / OpenAI / Ollama\ncompletion requests", PURPLE),
    ("Docker Daemon", "Build images, run\ncontainers, health checks", GREEN),
    ("GitHub", "Repository cloning\nfor project registration", ORANGE),
    ("Redis / RQ", "Async job queue\nfor deployments", RED),
]

tf = tb(s, Inches(0.6), Inches(1.4), Inches(3), Inches(0.4))
txt(tf, "Actors", sz=18, bold=True, color=ACCENT)

for i, (name, desc, col) in enumerate(actors_data):
    top = Inches(1.9) + i * Inches(1.0)
    rect(s, Inches(0.6), top, Inches(4.5), Inches(0.85), BG_CARD, border=col)
    tf = tb(s, Inches(0.8), top + Inches(0.05), Inches(1.5), Inches(0.35))
    txt(tf, name, sz=14, bold=True, color=col)
    tf2 = tb(s, Inches(2.3), top + Inches(0.05), Inches(2.6), Inches(0.75))
    tf2.word_wrap = True
    txt(tf2, desc, sz=11, color=LIGHT)

# Use Cases
usecases = [
    "Register Project (local or GitHub)",
    "Upload Environment File",
    "Parse Natural Language Command",
    "Preview Plan with AI Risk Score",
    "Approve Deployment Plan",
    "Execute Deployment Pipeline",
    "Monitor Deployment Health",
    "Trigger Manual Rollback",
    "Multi-Turn DevOps Chat",
]

tf = tb(s, Inches(5.5), Inches(1.4), Inches(3), Inches(0.4))
txt(tf, "Use Cases", sz=18, bold=True, color=GREEN)

for i, uc in enumerate(usecases):
    top = Inches(1.9) + i * Inches(0.56)
    rect(s, Inches(5.5), top, Inches(7), Inches(0.48), BG_CARD, border=GREEN)
    tf = tb(s, Inches(5.7), top + Inches(0.05), Inches(6.5), Inches(0.4))
    r = txt(tf, "", sz=13, color=LIGHT)
    run = r.add_run(); run.text = f"UC{i+1}:  "; run.font.size = Pt(13)
    run.font.color.rgb = GREEN; run.font.bold = True; run.font.name = "Segoe UI"
    run2 = r.add_run(); run2.text = uc; run2.font.size = Pt(13)
    run2.font.color.rgb = LIGHT; run2.font.name = "Segoe UI"

# Placeholder note
tf = tb(s, Inches(5.5), Inches(7.0), Inches(7), Inches(0.4))
txt(tf, "\u2191 See full Use Case Diagram in project report (Fig 1)", sz=11, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — System Design: Data Flow Diagram
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "07", "System Design \u2014 Data Flow Diagrams")

# DFD Level 0
rect(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD, border=ACCENT)
tf = tb(s, Inches(0.8), Inches(1.6), Inches(5.4), Inches(0.4))
txt(tf, "DFD Level 0 \u2014 Context Diagram", sz=16, bold=True, color=ACCENT)

l0_items = [
    "External Inputs:",
    "  \u2022 User commands (NL text, project payloads, approvals)",
    "  \u2022 LLM completion responses",
    "  \u2022 Docker daemon status",
    "",
    "System Outputs:",
    "  \u2022 Running containers",
    "  \u2022 Execution logs & rollback events",
    "  \u2022 Plan previews & risk assessments",
    "  \u2022 Chat responses with project context",
]
tf2 = tb(s, Inches(0.8), Inches(2.1), Inches(5.4), Inches(4.5))
tf2.word_wrap = True
for item in l0_items:
    if item == "":
        txt(tf2, " ", sz=8, color=BG_CARD, after=Pt(2))
    elif item.startswith("  "):
        txt(tf2, item, sz=13, color=LIGHT, after=Pt(3))
    else:
        txt(tf2, item, sz=14, bold=True, color=ORANGE, after=Pt(4))

# DFD Level 1
rect(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD, border=GREEN)
tf = tb(s, Inches(7.0), Inches(1.6), Inches(5.4), Inches(0.4))
txt(tf, "DFD Level 1 \u2014 Major Processes", sz=16, bold=True, color=GREEN)

processes = [
    ("P1", "Command Interpretation", "NL text \u2192 structured plan"),
    ("P2", "Plan Validation & Risk", "Evaluate against best practices"),
    ("P3", "Deployment Execution", "Docker pipeline on approval"),
    ("P4", "Health Monitoring", "Continuous eval + rollback"),
]

for i, (pid, name, desc) in enumerate(processes):
    top = Inches(2.2) + i * Inches(0.95)
    rect(s, Inches(7.0), top, Inches(5.4), Inches(0.8), RGBColor(0x15, 0x1F, 0x35), border=GREEN)
    badge = rect(s, Inches(7.1), top + Inches(0.15), Inches(0.5), Inches(0.5), GREEN)
    btf = badge.text_frame; bp = btf.paragraphs[0]; bp.text = pid
    bp.font.size = Pt(12); bp.font.bold = True; bp.font.color.rgb = BG_DARK
    bp.alignment = PP_ALIGN.CENTER; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = tb(s, Inches(7.7), top + Inches(0.08), Inches(4.5), Inches(0.35))
    txt(tf, name, sz=14, bold=True, color=WHITE)
    tf2 = tb(s, Inches(7.7), top + Inches(0.4), Inches(4.5), Inches(0.35))
    txt(tf2, desc, sz=12, color=GRAY)

# Data stores
tf3 = tb(s, Inches(7.0), Inches(6.1), Inches(5.4), Inches(0.6))
tf3.word_wrap = True
txt(tf3, "Data Stores:  Project DB (SQLite)  \u2022  Execution Logs  \u2022  Encrypted Env Files",
    sz=12, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Implementation: System Setup
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "08", "Implementation \u2014 System Setup & Configuration")

card(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(3.0),
     "Application Initialization (main.py)", [
         "FastAPI lifespan context manager for startup/shutdown",
         "Stale devops-* container cleanup on startup",
         "Persistent asyncio task for 60s watchdog loop",
         "CORS middleware + rate-limiting middleware",
         "SQLite database table initialization via init_db()",
     ], ac=ACCENT, isz=14)

card(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.0),
     "Safety Configuration (settings.py)", [
         "DRY_RUN = true by default (log-only execution)",
         "ENABLE_LOCAL_EXECUTION = false by default",
         "ALLOWED_REPO_ROOTS = empty (permit any in demo)",
         "Pydantic-settings reads from .env file",
         "Enforced at execution time, not configuration time",
     ], ac=RED, isz=14)

card(s, Inches(0.6), Inches(4.8), Inches(12), Inches(2.3),
     "Docker Compose \u2014 Full Stack (one-command deployment)", [
         "Redis:  redis:7-alpine with health check  |  Backend:  FastAPI + Docker socket mount (port 3001)",
         "Worker:  RQ background worker sharing backend-data volume  |  Frontend:  Next.js (port 3000)",
         "Single command:  docker-compose up --build  \u2014  starts entire stack with health dependencies",
     ], ac=GREEN, isz=15)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Implementation: Command Interpretation
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "08", "Implementation \u2014 LLM-Powered Command Interpretation")

# Flow
card(s, Inches(0.6), Inches(1.5), Inches(12), Inches(2.0),
     "Dual-Strategy Pipeline", [
         "1. User submits NL text \u2192  2. LLM parse via Gemini (system prompt + 6 few-shot examples) \u2192  3. Validate confidence score",
         "4. If confidence < 0.4 OR LLM fails \u2192  deterministic regex fallback  \u2192  5. Result tagged with interpretation_method",
     ], ac=PURPLE, isz=15)

# LLM output schema
card(s, Inches(0.6), Inches(3.8), Inches(5.8), Inches(3.3),
     "LLM Output Schema", [
         "action: deploy | rollback | status | ...",
         "version: semantic version or 'latest'",
         "environments: [staging, production, ...]",
         "post_steps: [run tests, smoke tests, ...]",
         "confidence: 0.0 \u2013 1.0 (threshold: 0.4)",
         "ai_reasoning: explanation of interpretation",
     ], ac=ACCENT, isz=14)

# Deterministic fallback
card(s, Inches(6.8), Inches(3.8), Inches(5.8), Inches(3.3),
     "Deterministic Regex Fallback", [
         "Activated when LLM unavailable or low confidence",
         "Regex patterns extract action, version, environments",
         "Result tagged: interpretation_method = 'deterministic'",
         "Frontend displays appropriate label to user",
         "Ensures system never fails on command parsing",
     ], ac=ORANGE, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Implementation: Risk Assessment
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "08", "Implementation \u2014 RAG-Style Plan Risk Assessment")

card(s, Inches(0.6), Inches(1.5), Inches(6), Inches(2.5),
     "LLM Risk Advisor", [
         "Receives JSON summary of deployment plan",
         "System prompt: 'reason as a senior DevOps engineer'",
         "Returns: risk_level, risk_score (0\u2013100),",
         "  warnings[], recommendations[], approval_recommendation",
     ], ac=ORANGE, isz=14)

card(s, Inches(7), Inches(1.5), Inches(5.6), Inches(2.5),
     "Deterministic Fallback Scoring", [
         "Production environment: +30 points",
         "Missing tests: +20 points",
         "No smoke tests (prod): +10 points",
         "Unknown action type: +25 points",
         "Missing version: +10 points",
     ], ac=RED, isz=14)

# Risk levels
rect(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.5), BG_CARD, border=ACCENT)
tf = tb(s, Inches(0.8), Inches(4.35), Inches(11.5), Inches(0.4))
r = txt(tf, "", sz=14, color=LIGHT)
for label, col, threshold in [("LOW", GREEN, "\u226430"), ("MEDIUM", ORANGE, "\u226460"),
                                ("HIGH", RED, "\u226480"), ("CRITICAL", RGBColor(0xFF, 0x00, 0x00), "\u2264100")]:
    run = r.add_run(); run.text = f"  \u25CF {label} ({threshold})  "; run.font.size = Pt(14)
    run.font.color.rgb = col; run.font.bold = True; run.font.name = "Segoe UI"

# Plan preview
card(s, Inches(0.6), Inches(5.1), Inches(12), Inches(2.0),
     "Plan Preview (returned to user before execution)", [
         "AI Confidence Score  \u2022  AI Reasoning  \u2022  Risk Level Badge  \u2022  Risk Score",
         "Warnings list  \u2022  Recommendations  \u2022  Explicit Approve / Cancel buttons",
         "No deployment executes without user clicking 'Approve' \u2014 mandatory safety gate",
     ], ac=GREEN, isz=15)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — Implementation: Orchestrator Pipeline
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "08", "Implementation \u2014 Six-Stage Deployment Pipeline")

stages = [
    ("1", "Validate", "Path within ALLOWED_REPO_ROOTS\n.env file present\nWorkspace resolution", GREEN),
    ("2", "Git Pull", "Applicable for GitHub\nprojects only\nSkipped for local", GRAY),
    ("3", "Docker Build", "Docker Python SDK\nAuto-generate Dockerfile\nif none exists", ACCENT),
    ("4", "Stop Old", "Retrieve latest deployment\nGracefully remove\nprevious container", ORANGE),
    ("5", "Docker Run", "Launch container with\nenv file + port mapping\n(EXPOSE \u2192 .env \u2192 default)", PURPLE),
    ("6", "Health Check", "Poll with 5 retries\nRegister deployment\nStore last-known-good tag", GREEN),
]

for i, (num, title, desc, col) in enumerate(stages):
    left = Inches(0.4) + i * Inches(2.1)
    rect(s, left, Inches(1.6), Inches(1.95), Inches(3.0), BG_CARD, border=col)
    badge = rect(s, left + Inches(0.65), Inches(1.7), Inches(0.55), Inches(0.55), col)
    btf = badge.text_frame; bp = btf.paragraphs[0]; bp.text = num
    bp.font.size = Pt(18); bp.font.bold = True; bp.font.color.rgb = BG_DARK
    bp.alignment = PP_ALIGN.CENTER; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = tb(s, left + Inches(0.05), Inches(2.35), Inches(1.85), Inches(0.35))
    txt(tf, title, sz=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    tf2 = tb(s, left + Inches(0.05), Inches(2.75), Inches(1.85), Inches(1.7))
    tf2.word_wrap = True
    txt(tf2, desc, sz=11, color=LIGHT, align=PP_ALIGN.CENTER)

# Port resolution + safety
card(s, Inches(0.6), Inches(4.9), Inches(5.8), Inches(2.2),
     "Port Resolution Priority", [
         "1. Dockerfile EXPOSE directive (authoritative)",
         "2. PORT= or APP_PORT= in project .env",
         "3. Stack default: 8080 (Node), 8000 (Python)",
         "Conflicting ports auto-cleared before launch",
     ], ac=ACCENT, isz=14)

card(s, Inches(6.8), Inches(4.9), Inches(5.8), Inches(2.2),
     "Safety Controls in Pipeline", [
         "DRY_RUN check: log actions, skip execution",
         "ENABLE_LOCAL_EXECUTION gate verified first",
         "Docker Python SDK only \u2014 no raw shell commands",
         "All operations use correlation IDs for tracing",
     ], ac=RED, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — Implementation: Watchdog & Rollback
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "08", "Implementation \u2014 Watchdog & Automatic Rollback")

card(s, Inches(0.6), Inches(1.5), Inches(6), Inches(3.5),
     "Watchdog Service (evaluate_deployments)", [
         "Runs every 60 seconds via asyncio in FastAPI lifespan",
         "Iterates all projects, checks latest 'running' deployment",
         "Calls container_health_check() via Docker SDK",
         "If unhealthy: marks deployment as 'failed'",
         "Calls _rollback_to_last_known_good() automatically",
         "Logs all actions via structlog with component binding",
     ], ac=RED, isz=14)

card(s, Inches(7), Inches(1.5), Inches(5.6), Inches(3.5),
     "Rollback Mechanism", [
         "Retrieves last_known_good_tag from project record",
         "Decrypts and mounts environment file",
         "Starts new container with known-good image",
         "Creates new Deployment record (status: 'running')",
         "Only rolls back if LKG tag differs from failed image",
         "Entire process is fully automated \u2014 no human action",
     ], ac=GREEN, isz=14)

# State machines
card(s, Inches(0.6), Inches(5.3), Inches(5.8), Inches(1.8),
     "Execution State Machine", [
         "queued \u2192 running | failed",
         "running \u2192 succeeded | failed | rolled_back",
         "succeeded \u2192 rolled_back    |    failed \u2192 rolled_back",
     ], ac=PURPLE, isz=14)

card(s, Inches(6.8), Inches(5.3), Inches(5.8), Inches(1.8),
     "Plan State Machine", [
         "pending_approval \u2192 approved | cancelled",
         "approved \u2192 running | cancelled",
         "running \u2192 succeeded | failed  (terminal states)",
     ], ac=ORANGE, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — Implementation: Chat Engine
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "08", "Implementation \u2014 Multi-Turn Chat Interface")

card(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.8),
     "Conversation Engine (Backend)", [
         "Maintains ChatSession + ChatMessage DB records",
         "Loads full history on each turn for context",
         "Injects real-time project state into system prompt:",
         "  current container ID, LKG tag, recent executions",
         "Parses <<<DEPLOY_INTENT>>>...<<<END_INTENT>>> markers",
         "Auto-creates Plan record when intent detected",
     ], ac=ACCENT, isz=14)

card(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.8),
     "Chat Frontend (React Components)", [
         "Real-time message rendering with typing indicators",
         "Deployment intent cards embedded in chat flow",
         "Plan preview + approval inline in conversation",
         "Session management: create / resume / list",
         "API calls centralized in lib/api.ts",
     ], ac=GREEN, isz=14)

card(s, Inches(0.6), Inches(4.6), Inches(12), Inches(2.5),
     "LLM Provider Abstraction Layer (llm/)", [
         "BaseLLMClient abstract class: complete() + complete_json() with 3-attempt JSON extraction",
         "Strategy: (1) direct JSON parse  \u2192  (2) markdown fence extraction  \u2192  (3) first {..} block  \u2192  (4) retry with strict prompt",
         "GeminiClient: google-generativeai SDK  |  OpenAIClient: raw httpx to any OpenAI-compatible endpoint",
         "OllamaClient: reuses OpenAIClient with overridden base_url and model  |  Factory selects via LLM_PROVIDER env var",
     ], ac=PURPLE, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 17 — Implementation: Screenshots placeholder
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "08", "Implementation \u2014 Application Screenshots")

placeholders = [
    ("Dashboard \u2014 Project List & Registration", Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.6)),
    ("Chat Interface \u2014 Multi-turn Conversation", Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.6)),
    ("Plan Preview \u2014 Risk Score & Approval Gate", Inches(0.6), Inches(4.4), Inches(5.8), Inches(2.6)),
    ("Execution Logs \u2014 6-Stage Pipeline Output", Inches(6.8), Inches(4.4), Inches(5.8), Inches(2.6)),
]

for title, l, t, w, h in placeholders:
    rect(s, l, t, w, h, BG_CARD, border=GRAY)
    tf = tb(s, l + Inches(0.2), t + h/2 - Inches(0.4), w - Inches(0.4), Inches(0.8))
    tf.word_wrap = True
    txt(tf, "[Screenshot]  " + title, sz=16, color=GRAY, align=PP_ALIGN.CENTER)
    txt(tf, "(Insert screenshot here)", sz=12, color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 18 — Testing (1/2)
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "08", "Testing \u2014 Unit & Integration Tests")

card(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(3.0),
     "Unit Tests (pytest + anyio)", [
         "test_demo_analysis.py: stack detection + Dockerfile gen",
         "Verifies Node.js \u2192 'node' stack, multi-stage Dockerfile",
         "test_m3_orchestrator.py: full Docker pipeline test",
         "Mocks DB & env, exercises real Docker operations",
         "Auto-skips when Docker daemon not accessible",
     ], ac=ACCENT, isz=14)

card(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.0),
     "Integration Tests", [
         "Command interpreter dual-strategy fallback",
         "LLM mocked to simulate failures \u2192 regex fallback",
         "Verifies interpretation_method = 'deterministic'",
         "Low confidence mocking \u2192 threshold logic validation",
         "Risk advisor LLM/deterministic interaction",
     ], ac=GREEN, isz=14)

card(s, Inches(0.6), Inches(4.8), Inches(12), Inches(2.3),
     "Functional (E2E) Test Scenarios", [
         "Happy path: register project \u2192 upload .env \u2192 NL command \u2192 review plan \u2192 approve \u2192 verify container",
         "Risk escalation: 'deploy to production no tests' \u2192 verify risk score > 60 and level = high/critical",
         "Fallback: invalid LLM_PROVIDER \u2192 deterministic parser handles gracefully",
         "Watchdog: manually stop container \u2192 wait 60s \u2192 verify auto-rollback to LKG image",
         "Chat intent: 'deploy latest to staging' in chat \u2192 verify Plan auto-created and linked",
     ], ac=ORANGE, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 19 — Testing (2/2): Performance & Security
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "08", "Testing \u2014 Performance & Security")

card(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.8),
     "Performance Results", [
         "LLM command parse: 3\u20136 seconds (Gemini API)",
         "Deterministic fallback: < 50 milliseconds",
         "RQ job enqueue: < 100 milliseconds",
         "Cached Docker build: < 10 seconds (Node.js)",
     ], ac=ACCENT, isz=15)

card(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.8),
     "Security Testing", [
         "Dry-run: no containers created with DRY_RUN=true",
         "Execution gate: fails cleanly when EXEC=false",
         "Path traversal: rejects paths outside whitelist",
         "Rate limiting: HTTP 429 after 30 req/min/IP",
         "Secrets: never logged via structlog output",
     ], ac=RED, isz=15)

card(s, Inches(0.6), Inches(4.6), Inches(12), Inches(2.5),
     "User Acceptance Testing", [
         "Scenario: developer with no Docker experience deploys a Node.js app to staging",
         "Types 'ship the app to staging and run tests' in chat \u2192 reviews AI plan \u2192 approves \u2192 observes live logs",
         "Follow-up: 'what was the last deployment status?' \u2192 receives contextual NL response from real project state",
         "Feedback: NL interface significantly lowered perceived complexity of deployment process",
     ], ac=GREEN, isz=15)


# ════════════════════════════════════════════════════════════════
# SLIDE 20 — Conclusion
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "09", "Conclusion")

card(s, Inches(0.6), Inches(1.5), Inches(12), Inches(2.3),
     "What We Achieved", [
         "Production-grade deployment automation made accessible AND safe through LLM agents + deterministic fallbacks",
         "Successfully bridges the gap between natural language intent and complex infrastructure operations",
         "Dual-strategy architecture ensures reliability: LLM primary, deterministic fallback \u2014 system never fails",
     ], ac=GREEN, isz=16)

achievements = [
    ("Intelligent\nUnderstanding", "LLM interprets intent;\nnot just execute commands\nbut reasons about them", ACCENT),
    ("Risk\nAwareness", "RAG-style advisor adds\ndeployment intelligence\nabsent from CI/CD tools", ORANGE),
    ("Mandatory\nSafety Gates", "Approval + dry-run +\nexecution gate prevent\naccidental destruction", RED),
    ("Self-Healing", "Watchdog auto-rollback\nensures temporary failures\nrecover automatically", GREEN),
    ("Conversational\nAssistant", "Multi-turn chat with\nreal-time project state\nawareness", PURPLE),
]

for i, (title, desc, col) in enumerate(achievements):
    left = Inches(0.4) + i * Inches(2.55)
    rect(s, left, Inches(4.2), Inches(2.35), Inches(2.6), BG_CARD, border=col)
    tf = tb(s, left + Inches(0.1), Inches(4.3), Inches(2.15), Inches(0.7))
    tf.word_wrap = True
    txt(tf, title, sz=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    tf2 = tb(s, left + Inches(0.1), Inches(5.0), Inches(2.15), Inches(1.6))
    tf2.word_wrap = True
    txt(tf2, desc, sz=12, color=LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 21 — Future Enhancements
# ════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "10", "Future Enhancements")

phases = [
    ("Cloud Provider\nIntegration", ACCENT, [
        "Native Vercel & Render deployment APIs",
        "Unified deploy status endpoint",
        "Cross-platform deployment orchestration",
        "Normalized status schema across providers",
    ]),
    ("Production\nHardening", GREEN, [
        "PostgreSQL migration (Alembic)",
        "WebSocket real-time log streaming",
        "Horizontal RQ worker scaling",
        "Git branch-aware deployments",
    ]),
    ("Observability &\nAnalytics", ORANGE, [
        "Deployment history timeline dashboard",
        "Prometheus metrics export + Grafana",
        "Slack / email alerting on rollback",
        "Success rate & MTTR visualization",
    ]),
    ("Advanced AI\nCapabilities", PURPLE, [
        "Multi-project dependency reasoning",
        "SOP ingestion via Chroma vector DB",
        "Predictive risk scoring from history",
        "True RAG with deployment SOPs",
    ]),
]

for i, (title, color, items) in enumerate(phases):
    left = Inches(0.5) + i * Inches(3.15)
    rect(s, left, Inches(1.6), Inches(2.95), Inches(5.2), BG_CARD, border=color)
    tf = tb(s, left + Inches(0.15), Inches(1.75), Inches(2.65), Inches(0.8))
    tf.word_wrap = True
    txt(tf, title, sz=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    line(s, left + Inches(0.3), Inches(2.55), Inches(2.35), color=color, thick=Pt(1.5))
    bullets(s, items, left + Inches(0.15), Inches(2.75), Inches(2.65), Inches(3.8),
            sz=13, col=LIGHT, bcol=GRAY, sp=Pt(8))


# ════════════════════════════════════════════════════════════════
# SLIDE 22 — Thank You
# ════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT)

rect(s, Inches(2.5), Inches(1.8), Inches(8.3), Inches(4.2), BG_CARD,
     border=RGBColor(0x2A, 0x3A, 0x5C))

tf = tb(s, Inches(3), Inches(2.2), Inches(7.3), Inches(1.0))
txt(tf, "Thank You", sz=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

line(s, Inches(5), Inches(3.2), Inches(3.3), thick=Pt(2))

tf2 = tb(s, Inches(3), Inches(3.5), Inches(7.3), Inches(2.0))
tf2.word_wrap = True
txt(tf2, "AI DevOps Commander", sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, after=Pt(8))
txt(tf2, "Natural Language Driven Deployment Automation with Agentic AI", sz=16, color=LIGHT, align=PP_ALIGN.CENTER, after=Pt(20))
txt(tf2, "Prince Parmar Singh  (24100443)", sz=16, color=GRAY, align=PP_ALIGN.CENTER, after=Pt(4))
txt(tf2, "Guide: Dr. Deepika MP, Assistant Professor", sz=14, color=GRAY, align=PP_ALIGN.CENTER, after=Pt(4))
txt(tf2, "Department of Computer Applications, CUSAT", sz=14, color=GRAY, align=PP_ALIGN.CENTER)

# Question prompt
tf3 = tb(s, Inches(3), Inches(6.3), Inches(7.3), Inches(0.5))
txt(tf3, "Questions & Discussion", sz=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ACCENT)


# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
out = "AI_DevOps_Commander_Presentation.pptx"
prs.save(out)
print(f"Presentation saved: {out}")
print(f"Total slides: {len(prs.slides)}")
