"""
AI DevOps Commander — Cream Neutral Minimalist Presentation
Inspired by Canva 'Cream Neutral Minimalist New Business Pitch Deck'
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Cream Neutral Minimalist Palette ────────────────────────────
CREAM      = RGBColor(0xF5, 0xF0, 0xE8)   # warm cream background
CREAM_DARK = RGBColor(0xEA, 0xE0, 0xD2)   # slightly darker cream for cards
BEIGE      = RGBColor(0xD4, 0xC5, 0xAF)   # beige accent / borders
TAN        = RGBColor(0xC2, 0xAE, 0x94)   # tan for decorative elements
BROWN      = RGBColor(0x6B, 0x5B, 0x47)   # warm brown for headings
DARK_BROWN = RGBColor(0x3E, 0x34, 0x28)   # near-black brown for body text
CHARCOAL   = RGBColor(0x2C, 0x25, 0x1E)   # darkest text
OLIVE      = RGBColor(0x7A, 0x8B, 0x6A)   # muted olive green accent
RUST       = RGBColor(0xB8, 0x73, 0x52)   # warm rust/terracotta accent
SAGE       = RGBColor(0xA3, 0xB1, 0x8A)   # sage green
DUSTY_ROSE = RGBColor(0xC4, 0x8B, 0x7F)   # dusty rose
SLATE      = RGBColor(0x8A, 0x83, 0x78)   # warm gray for subtle text
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BLACK = RGBColor(0x33, 0x2D, 0x26)   # soft black

SERIF = "Georgia"
SANS  = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


# ── Helpers ─────────────────────────────────────────────────────
def set_bg(slide, color=CREAM):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill, border=None, bw=Pt(1), radius=0.04):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.fill.background()
    s.adjustments[0] = radius
    return s


def add_sharp_rect(slide, l, t, w, h, fill, border=None, bw=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.fill.background()
    return s


def thin_line(slide, l, t, w, color=TAN, thick=Pt(1.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, thick)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def circle(slide, l, t, size, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def tfb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h).text_frame


def txt(tf, text, sz=16, color=DARK_BROWN, bold=False, italic=False,
        align=PP_ALIGN.LEFT, after=Pt(6), before=Pt(0), font=SANS):
    if tf.paragraphs[0].text == "" and not hasattr(tf, '_u'):
        p = tf.paragraphs[0]; tf._u = True
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    p.space_after = after
    p.space_before = before
    return p


def slide_header(slide, num, title, subtitle=None):
    """Minimalist header: small number, serif title, thin underline."""
    set_bg(slide)
    # Number
    tf = tfb(slide, Inches(0.8), Inches(0.4), Inches(0.6), Inches(0.5))
    txt(tf, num, sz=14, color=SLATE, font=SANS, bold=False)
    # Title
    tf2 = tfb(slide, Inches(1.3), Inches(0.3), Inches(10), Inches(0.65))
    txt(tf2, title, sz=32, color=CHARCOAL, bold=False, font=SERIF)
    # Thin line
    thin_line(slide, Inches(0.8), Inches(1.0), Inches(11.5), color=BEIGE, thick=Pt(1))
    if subtitle:
        tf3 = tfb(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4))
        txt(tf3, subtitle, sz=14, color=SLATE, italic=True)


def bullet_list(slide, items, l, t, w, h, sz=15, col=DARK_BROWN, bcol=TAN, sp=Pt(8)):
    tf = tfb(slide, l, t, w, h)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = sp
        p.space_before = Pt(2)
        r1 = p.add_run()
        r1.text = "\u2014  "  # em dash bullet (minimalist)
        r1.font.size = Pt(sz)
        r1.font.color.rgb = bcol
        r1.font.name = SANS
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(sz)
        r2.font.color.rgb = col
        r2.font.name = SANS


def card(slide, l, t, w, h, title, items, ac=BROWN, tsz=16, isz=14):
    add_rect(slide, l, t, w, h, CREAM_DARK, border=BEIGE)
    # Accent strip on left
    add_sharp_rect(slide, l, t + Inches(0.1), Pt(4), h - Inches(0.2), ac)
    tf = tfb(slide, l + Inches(0.25), t + Inches(0.15), w - Inches(0.4), Inches(0.4))
    txt(tf, title, sz=tsz, bold=True, color=ac, font=SERIF)
    bullet_list(slide, items, l + Inches(0.25), t + Inches(0.55),
                w - Inches(0.45), h - Inches(0.7), sz=isz, col=DARK_BROWN, bcol=TAN, sp=Pt(5))


def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


# Decorative corner dots (minimalist touch)
def corner_dots(slide):
    """Add subtle decorative dots in corners."""
    for cx, cy in [(Inches(0.3), Inches(0.3)),
                   (Inches(12.8), Inches(0.3)),
                   (Inches(0.3), Inches(7.0)),
                   (Inches(12.8), Inches(7.0))]:
        circle(slide, cx, cy, Inches(0.12), BEIGE)


def page_number(slide, num):
    tf = tfb(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.35))
    txt(tf, str(num), sz=11, color=SLATE, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s)

# Large decorative circle (top right)
circle(s, Inches(9.5), Inches(-1.5), Inches(5), CREAM_DARK, border=BEIGE)
# Smaller accent circle
circle(s, Inches(0.5), Inches(5.5), Inches(2.5), CREAM_DARK, border=BEIGE)

# Vertical thin line
add_sharp_rect(s, Inches(1.2), Inches(1.5), Pt(1.5), Inches(4.5), TAN)

# Title block
tf = tfb(s, Inches(1.8), Inches(1.5), Inches(8), Inches(1.5))
tf.word_wrap = True
txt(tf, "AI DevOps Commander", sz=48, color=CHARCOAL, font=SERIF, after=Pt(4))
txt(tf, "Natural Language Driven Deployment", sz=24, color=BROWN, font=SERIF, italic=True, after=Pt(2))
txt(tf, "Automation with Agentic AI", sz=24, color=BROWN, font=SERIF, italic=True)

# Thin separator
thin_line(s, Inches(1.8), Inches(3.65), Inches(3), color=TAN, thick=Pt(1))

# Submitted by
tf2 = tfb(s, Inches(1.8), Inches(3.9), Inches(8), Inches(1.8))
tf2.word_wrap = True
txt(tf2, "Submitted by", sz=12, color=SLATE, after=Pt(2))
txt(tf2, "Prince Parmar Singh", sz=20, color=CHARCOAL, bold=True, font=SERIF, after=Pt(2))
txt(tf2, "Roll No: 24100443", sz=14, color=SLATE, after=Pt(16))
txt(tf2, "Under the guidance of", sz=12, color=SLATE, after=Pt(2))
txt(tf2, "Dr. Deepika MP, Assistant Professor", sz=16, color=CHARCOAL, font=SERIF)

# Footer
tf3 = tfb(s, Inches(1.8), Inches(6.4), Inches(8), Inches(0.6))
tf3.word_wrap = True
txt(tf3, "Department of Computer Applications", sz=12, color=SLATE, after=Pt(2))
txt(tf3, "Cochin University of Science and Technology  |  April 2026", sz=12, color=SLATE)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Introduction
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "01", "Introduction")
corner_dots(s)
page_number(s, 2)

# Left — the problem context
card(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.5),
     "The Deployment Challenge", [
         "CI/CD demands deep CLI, YAML, and cloud expertise",
         "Deployment errors cost teams hours of debugging",
         "Non-DevOps devs are locked out of workflows",
         "Monitoring and rollback are fragmented",
     ], ac=RUST)

# Right — the opportunity
card(s, Inches(6.8), Inches(1.5), Inches(5.6), Inches(2.5),
     "The LLM Opportunity", [
         "LLMs can reason about structured plans from text",
         "Conversational AI bridges intent and infrastructure",
         "Agentic pipelines enable safe automation",
         "Natural language lowers the barrier for everyone",
     ], ac=OLIVE)

# Bottom — what it is
add_rect(s, Inches(0.8), Inches(4.3), Inches(11.6), Inches(2.8), CREAM_DARK, border=BEIGE)
add_sharp_rect(s, Inches(0.8), Inches(4.3), Pt(4), Inches(2.8), BROWN)
tf = tfb(s, Inches(1.1), Inches(4.4), Inches(11), Inches(0.4))
txt(tf, "AI DevOps Commander", sz=18, bold=True, color=BROWN, font=SERIF)
bullet_list(s, [
    "Full-stack agentic system: deploy software using plain English commands",
    "Next.js 14 dashboard  +  FastAPI backend  +  Redis job queue  +  Docker SDK",
    "Multi-strategy: LLM primary (Google Gemini) with deterministic regex fallback",
    "LLM risk assessment, mandatory approval gates, automatic watchdog rollback",
    "Layered safety: dry-run default, execution gates, path whitelists, encrypted secrets",
], Inches(1.1), Inches(4.9), Inches(11), Inches(2.0), sz=14, col=DARK_BROWN, bcol=TAN, sp=Pt(5))


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "02", "Problem Statement")
corner_dots(s)
page_number(s, 3)

problems = [
    ("Technical Barriers", "Deployment requires Docker CLI, CI/CD YAML, cloud APIs, and shell scripting \u2014 creating an access gap.", RUST),
    ("Cognitive Overload", "Managing environments, tracking versions, and coordinating rollbacks places unsustainable demands.", BROWN),
    ("No Contextual Intelligence", "Tools execute blindly \u2014 deploying to production on a Friday without tests is allowed.", DUSTY_ROSE),
    ("Fragmented Observability", "Status, health, rollbacks, and logs are scattered across multiple dashboards.", OLIVE),
    ("No Conversational Interface", "No deployment tool supports natural language interaction for infrastructure.", SLATE),
]

for i, (title, desc, col) in enumerate(problems):
    row = i // 3
    coli = i % 3
    left = Inches(0.8) + coli * Inches(4.05)
    top = Inches(1.5) + row * Inches(2.7)
    w, h = Inches(3.8), Inches(2.4)
    add_rect(s, left, top, w, h, CREAM_DARK, border=BEIGE)
    # Accent circle with number
    c = circle(s, left + Inches(0.2), top + Inches(0.2), Inches(0.45), col)
    ctf = c.text_frame
    cp = ctf.paragraphs[0]
    cp.text = str(i + 1)
    cp.font.size = Pt(16)
    cp.font.bold = True
    cp.font.color.rgb = WHITE
    cp.font.name = SANS
    cp.alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Title
    tf = tfb(s, left + Inches(0.8), top + Inches(0.2), w - Inches(1), Inches(0.4))
    txt(tf, title, sz=15, bold=True, color=col, font=SERIF)
    # Desc
    tf2 = tfb(s, left + Inches(0.2), top + Inches(0.8), w - Inches(0.4), h - Inches(1))
    tf2.word_wrap = True
    txt(tf2, desc, sz=13, color=DARK_BROWN)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Objectives
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "03", "Objectives")
corner_dots(s)
page_number(s, 4)

pillars = [
    ("Natural Language\nInterface", RUST, [
        "Plain English deployment commands",
        "Multi-turn chat with intent detection",
        "Zero CLI knowledge required",
    ]),
    ("Intelligent\nRisk Assessment", OLIVE, [
        "LLM plan advisor scores risk 0\u2013100",
        "Evaluates env, tests, version, timing",
        "Deterministic fallback ensures reliability",
    ]),
    ("Safe Automated\nExecution", BROWN, [
        "6-stage Docker pipeline (Python SDK)",
        "Mandatory approval before execution",
        "Dry-run + execution gate + path whitelist",
    ]),
    ("Self-Healing\nInfrastructure", DUSTY_ROSE, [
        "60-second watchdog health polling",
        "Auto-rollback to last-known-good image",
        "Encrypted secrets + structured logging",
    ]),
]

for i, (title, color, items) in enumerate(pillars):
    left = Inches(0.6) + i * Inches(3.15)
    add_rect(s, left, Inches(1.5), Inches(2.95), Inches(5.3), CREAM_DARK, border=BEIGE)
    # Top accent bar
    add_sharp_rect(s, left + Inches(0.3), Inches(1.5), Inches(2.35), Pt(4), color)
    # Number
    tf = tfb(s, left + Inches(0.2), Inches(1.7), Inches(2.55), Inches(0.4))
    txt(tf, f"0{i+1}", sz=36, color=color, bold=False, font=SERIF, align=PP_ALIGN.CENTER)
    # Title
    tf2 = tfb(s, left + Inches(0.15), Inches(2.4), Inches(2.65), Inches(0.8))
    tf2.word_wrap = True
    txt(tf2, title, sz=15, bold=True, color=color, font=SERIF, align=PP_ALIGN.CENTER)
    # Thin line
    thin_line(s, left + Inches(0.5), Inches(3.25), Inches(1.95), color=BEIGE)
    # Items
    bullet_list(s, items, left + Inches(0.15), Inches(3.45), Inches(2.65), Inches(3),
                sz=13, col=DARK_BROWN, bcol=TAN, sp=Pt(8))


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Proposed System: Request Flow
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "04", "Proposed System")
corner_dots(s)
page_number(s, 5)

# Subtitle
tf = tfb(s, Inches(0.8), Inches(1.15), Inches(10), Inches(0.4))
txt(tf, "End-to-end request flow from natural language to running container", sz=14, color=SLATE, italic=True)

steps = [
    ("User Input", "Plain English\ncommand via chat", RUST),
    ("Interpreter", "LLM parse with\nregex fallback", BROWN),
    ("Risk Advisor", "RAG-style risk\nscoring (0\u2013100)", OLIVE),
    ("Approval", "Human review\n& approve", SAGE),
    ("Orchestrator", "6-stage Docker\npipeline", DUSTY_ROSE),
    ("Watchdog", "Health monitor\n& auto-rollback", SLATE),
]

bw = Inches(1.7)
bh = Inches(1.7)
gap = Inches(0.28)
sx = Inches(0.5)
y = Inches(1.8)

for i, (lbl, desc, col) in enumerate(steps):
    x = sx + i * (bw + gap)
    add_rect(s, x, y, bw, bh, CREAM_DARK, border=BEIGE)
    # Top accent strip
    add_sharp_rect(s, x, y, bw, Pt(3.5), col)
    # Step number
    tf = tfb(s, x + Inches(0.05), y + Inches(0.12), Inches(0.4), Inches(0.3))
    txt(tf, str(i + 1), sz=20, color=col, font=SERIF, bold=False)
    # Label
    tf2 = tfb(s, x + Inches(0.35), y + Inches(0.12), bw - Inches(0.4), Inches(0.3))
    txt(tf2, lbl, sz=12, bold=True, color=col, font=SERIF)
    # Desc
    tf3 = tfb(s, x + Inches(0.1), y + Inches(0.55), bw - Inches(0.2), Inches(1))
    tf3.word_wrap = True
    txt(tf3, desc, sz=11, color=DARK_BROWN)
    # Arrow
    if i < len(steps) - 1:
        ax = x + bw + Inches(0.02)
        tf_a = tfb(s, ax, y + Inches(0.6), gap - Inches(0.04), Inches(0.35))
        txt(tf_a, "\u2192", sz=18, color=TAN, align=PP_ALIGN.CENTER)

# Architecture layers (bottom)
layers = [
    ("Frontend", "Next.js 14  |  React 18  |  Tailwind CSS  |  TypeScript", RUST),
    ("Backend", "FastAPI  |  SQLModel  |  Pydantic v2  |  structlog  |  RQ + Redis", OLIVE),
    ("Infrastructure", "Docker Python SDK  |  SQLite  |  Encrypted Secrets  |  Watchdog", BROWN),
    ("LLM Layer", "Google Gemini (primary)  |  OpenAI  |  Ollama  |  Provider Abstraction", DUSTY_ROSE),
]

for i, (lbl, desc, col) in enumerate(layers):
    top = Inches(3.95) + i * Inches(0.82)
    add_rect(s, Inches(0.8), top, Inches(11.6), Inches(0.7), CREAM_DARK, border=BEIGE)
    add_sharp_rect(s, Inches(0.8), top, Pt(4), Inches(0.7), col)
    tf = tfb(s, Inches(1.1), top + Inches(0.12), Inches(2), Inches(0.4))
    txt(tf, lbl, sz=13, bold=True, color=col, font=SERIF)
    tf2 = tfb(s, Inches(3.2), top + Inches(0.12), Inches(9), Inches(0.4))
    txt(tf2, desc, sz=13, color=DARK_BROWN)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Feasibility Analysis
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "05", "Feasibility Analysis")
corner_dots(s)
page_number(s, 6)

feasibility = [
    ("Technical", OLIVE, [
        "Mature open-source stack: FastAPI, Next.js, Docker SDK, Redis",
        "Google Gemini API with structured output support",
        "Python 3.11+, Node.js 18+, Docker \u2014 widely available",
        "All frameworks are production-grade with active communities",
    ]),
    ("Operational", RUST, [
        "NL interface lowers barrier for non-DevOps users",
        "Risk advisor provides intelligent guardrails",
        "Watchdog provides automated recovery",
        "Validated via unit, integration, and E2E tests",
    ]),
    ("Economic", BROWN, [
        "Entirely open-source \u2014 zero license costs",
        "Minimal LLM API costs for dev/demo use",
        "Only needs: machine with Docker installed",
        "Full Docker Compose for one-command deploy",
    ]),
]

for i, (title, color, items) in enumerate(feasibility):
    left = Inches(0.6) + i * Inches(4.15)
    card(s, left, Inches(1.5), Inches(3.95), Inches(5.2), title + " Feasibility", items,
         ac=color, tsz=17, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Technologies Used
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "06", "Technologies Used")
corner_dots(s)
page_number(s, 7)

card(s, Inches(0.8), Inches(1.4), Inches(3.7), Inches(5.5),
     "Frontend", [
         "Next.js 14 \u2014 React framework",
         "React 18 \u2014 Component library",
         "TypeScript 5.5 \u2014 Type safety",
         "Tailwind CSS 3.4 \u2014 Styling",
         "Vitest 2.0 \u2014 Unit testing",
         "lucide-react \u2014 Icons",
     ], ac=RUST, isz=15)

card(s, Inches(4.8), Inches(1.4), Inches(3.7), Inches(5.5),
     "Backend", [
         "Python 3.12 \u2014 Runtime",
         "FastAPI 0.115 \u2014 API framework",
         "SQLModel \u2014 ORM (SQLAlchemy)",
         "Pydantic v2 \u2014 Validation",
         "RQ + Redis \u2014 Job queue",
         "structlog \u2014 Logging",
         "pytest + anyio \u2014 Testing",
     ], ac=OLIVE, isz=15)

card(s, Inches(8.8), Inches(1.4), Inches(3.7), Inches(5.5),
     "Infrastructure & AI", [
         "Docker Python SDK \u2014 Containers",
         "Redis 7 \u2014 Queue backend",
         "SQLite \u2014 Dev database",
         "cryptography \u2014 Encryption",
         "Google Gemini \u2014 Primary LLM",
         "OpenAI / Ollama \u2014 Alternates",
         "httpx \u2014 Async HTTP client",
     ], ac=BROWN, isz=15)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — System Design: Architecture
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "07", "System Design \u2014 Service Architecture")
corner_dots(s)
page_number(s, 8)

services = [
    ("command_interpreter.py", "Parses NL commands; LLM primary, regex fallback", RUST),
    ("orchestrator.py", "Executes approved plans via 6-stage Docker pipeline", OLIVE),
    ("rag_advisor.py", "LLM risk assessment with deterministic fallback", BROWN),
    ("watchdog.py", "Container health polling; auto-rollback on failure", DUSTY_ROSE),
    ("project_analysis.py", "Tech stack detection; auto-generates Dockerfiles", SAGE),
    ("conversation_engine.py", "Multi-turn chat with deployment intent detection", RUST),
    ("llm/factory.py", "Provider-agnostic LLM client selection factory", SLATE),
]

# Table header
add_rect(s, Inches(0.8), Inches(1.4), Inches(4.2), Inches(0.5), TAN)
tf = tfb(s, Inches(0.9), Inches(1.42), Inches(4), Inches(0.45))
txt(tf, "Service Module", sz=14, bold=True, color=WHITE, font=SERIF, align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.0), Inches(1.4), Inches(7.2), Inches(0.5), TAN)
tf = tfb(s, Inches(5.1), Inches(1.42), Inches(7), Inches(0.45))
txt(tf, "Responsibility", sz=14, bold=True, color=WHITE, font=SERIF, align=PP_ALIGN.CENTER)

for i, (mod, role, col) in enumerate(services):
    top = Inches(2.0) + i * Inches(0.58)
    bg_c = CREAM_DARK if i % 2 == 0 else CREAM
    add_rect(s, Inches(0.8), top, Inches(4.2), Inches(0.5), bg_c, border=BEIGE)
    add_sharp_rect(s, Inches(0.8), top, Pt(3.5), Inches(0.5), col)
    tf = tfb(s, Inches(1.05), top + Inches(0.06), Inches(3.9), Inches(0.4))
    txt(tf, mod, sz=13, bold=True, color=col, font=SANS)
    add_rect(s, Inches(5.0), top, Inches(7.2), Inches(0.5), bg_c, border=BEIGE)
    tf2 = tfb(s, Inches(5.2), top + Inches(0.06), Inches(6.9), Inches(0.4))
    txt(tf2, role, sz=13, color=DARK_BROWN)

# API Routes bar
add_rect(s, Inches(0.8), Inches(6.2), Inches(11.4), Inches(0.8), CREAM_DARK, border=BEIGE)
add_sharp_rect(s, Inches(0.8), Inches(6.2), Pt(4), Inches(0.8), OLIVE)
tf = tfb(s, Inches(1.1), Inches(6.25), Inches(2), Inches(0.3))
txt(tf, "API Routes", sz=14, bold=True, color=OLIVE, font=SERIF)
tf2 = tfb(s, Inches(1.1), Inches(6.55), Inches(10.8), Inches(0.35))
txt(tf2, "projects  \u2022  commands  \u2022  executions  \u2022  providers  \u2022  deploy_status  \u2022  demo  \u2022  chat",
    sz=13, color=DARK_BROWN)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Use Case Diagram
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "07", "System Design \u2014 Use Case Diagram")
corner_dots(s)
page_number(s, 9)

# Actors
actors = [
    ("User", "Developer / Operator via web dashboard", RUST),
    ("LLM Provider", "Gemini / OpenAI / Ollama APIs", OLIVE),
    ("Docker Daemon", "Build, run, health check containers", BROWN),
    ("GitHub", "Repository cloning for registration", DUSTY_ROSE),
    ("Redis / RQ", "Async job queue for deployments", SLATE),
]

tf = tfb(s, Inches(0.8), Inches(1.3), Inches(2), Inches(0.4))
txt(tf, "Actors", sz=18, bold=True, color=BROWN, font=SERIF)

for i, (name, desc, col) in enumerate(actors):
    top = Inches(1.8) + i * Inches(1.0)
    add_rect(s, Inches(0.8), top, Inches(4.8), Inches(0.85), CREAM_DARK, border=BEIGE)
    c = circle(s, Inches(1.0), top + Inches(0.18), Inches(0.5), col)
    ctf = c.text_frame; cp = ctf.paragraphs[0]
    cp.text = name[0]; cp.font.size = Pt(16); cp.font.bold = True
    cp.font.color.rgb = WHITE; cp.alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = tfb(s, Inches(1.65), top + Inches(0.08), Inches(3.8), Inches(0.35))
    txt(tf, name, sz=14, bold=True, color=col, font=SERIF)
    tf2 = tfb(s, Inches(1.65), top + Inches(0.42), Inches(3.8), Inches(0.35))
    txt(tf2, desc, sz=11, color=SLATE)

# Use cases
usecases = [
    "Register Project (local or GitHub)",
    "Upload Environment File",
    "Parse Natural Language Command",
    "Preview Plan with AI Risk Score",
    "Approve Deployment Plan",
    "Execute Deployment Pipeline",
    "Monitor Deployment Health",
    "Trigger Manual Rollback",
    "Multi-Turn DevOps Chat",
]

tf = tfb(s, Inches(6.2), Inches(1.3), Inches(3), Inches(0.4))
txt(tf, "Use Cases", sz=18, bold=True, color=OLIVE, font=SERIF)

for i, uc in enumerate(usecases):
    top = Inches(1.8) + i * Inches(0.58)
    add_rect(s, Inches(6.2), top, Inches(6.3), Inches(0.5), CREAM_DARK, border=BEIGE)
    add_sharp_rect(s, Inches(6.2), top, Pt(3.5), Inches(0.5), OLIVE)
    tf = tfb(s, Inches(6.45), top + Inches(0.06), Inches(5.9), Inches(0.4))
    r = txt(tf, "", sz=13, color=DARK_BROWN)
    run1 = r.add_run(); run1.text = f"UC{i+1}  "; run1.font.size = Pt(12)
    run1.font.color.rgb = OLIVE; run1.font.bold = True; run1.font.name = SANS
    run2 = r.add_run(); run2.text = uc; run2.font.size = Pt(13)
    run2.font.color.rgb = DARK_BROWN; run2.font.name = SANS


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Data Flow Diagram
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "07", "System Design \u2014 Data Flow Diagrams")
corner_dots(s)
page_number(s, 10)

# DFD Level 0
add_rect(s, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.3), CREAM_DARK, border=BEIGE)
add_sharp_rect(s, Inches(0.8), Inches(1.4), Inches(5.6), Pt(4), RUST)
tf = tfb(s, Inches(1.0), Inches(1.55), Inches(5.2), Inches(0.4))
txt(tf, "DFD Level 0 \u2014 Context Diagram", sz=16, bold=True, color=RUST, font=SERIF)

l0 = [
    ("Inputs:", True, BROWN),
    ("User commands (NL text, project payloads, approvals)", False, DARK_BROWN),
    ("LLM completion responses", False, DARK_BROWN),
    ("Docker daemon status updates", False, DARK_BROWN),
    ("", False, CREAM_DARK),
    ("Outputs:", True, BROWN),
    ("Running containers + execution logs", False, DARK_BROWN),
    ("Plan previews + risk assessments", False, DARK_BROWN),
    ("Chat responses with project context", False, DARK_BROWN),
    ("Rollback events + status updates", False, DARK_BROWN),
]

tf2 = tfb(s, Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.5))
tf2.word_wrap = True
for text, is_header, col in l0:
    if text == "":
        txt(tf2, " ", sz=6, color=CREAM_DARK, after=Pt(4))
    elif is_header:
        txt(tf2, text, sz=14, bold=True, color=col, font=SERIF, after=Pt(4))
    else:
        p = tf2.add_paragraph()
        p.space_after = Pt(3)
        r1 = p.add_run(); r1.text = "\u2014  "; r1.font.size = Pt(13)
        r1.font.color.rgb = TAN; r1.font.name = SANS
        r2 = p.add_run(); r2.text = text; r2.font.size = Pt(13)
        r2.font.color.rgb = col; r2.font.name = SANS

# DFD Level 1
add_rect(s, Inches(6.8), Inches(1.4), Inches(5.6), Inches(5.3), CREAM_DARK, border=BEIGE)
add_sharp_rect(s, Inches(6.8), Inches(1.4), Inches(5.6), Pt(4), OLIVE)
tf = tfb(s, Inches(7.0), Inches(1.55), Inches(5.2), Inches(0.4))
txt(tf, "DFD Level 1 \u2014 Major Processes", sz=16, bold=True, color=OLIVE, font=SERIF)

processes = [
    ("P1", "Command Interpretation", "NL text \u2192 structured plan"),
    ("P2", "Plan Validation & Risk", "Evaluate against best practices"),
    ("P3", "Deployment Execution", "Docker pipeline on approval"),
    ("P4", "Health Monitoring", "Continuous eval + rollback"),
]

for i, (pid, name, desc) in enumerate(processes):
    top = Inches(2.2) + i * Inches(1.0)
    add_rect(s, Inches(7.0), top, Inches(5.2), Inches(0.85), CREAM, border=BEIGE)
    c = circle(s, Inches(7.15), top + Inches(0.17), Inches(0.5), OLIVE)
    ctf = c.text_frame; cp = ctf.paragraphs[0]; cp.text = pid
    cp.font.size = Pt(12); cp.font.bold = True; cp.font.color.rgb = WHITE
    cp.alignment = PP_ALIGN.CENTER; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = tfb(s, Inches(7.8), top + Inches(0.08), Inches(4.3), Inches(0.35))
    txt(tf, name, sz=14, bold=True, color=CHARCOAL, font=SERIF)
    tf2 = tfb(s, Inches(7.8), top + Inches(0.42), Inches(4.3), Inches(0.35))
    txt(tf2, desc, sz=12, color=SLATE)

# Data stores
tf3 = tfb(s, Inches(7.0), Inches(6.3), Inches(5.2), Inches(0.35))
txt(tf3, "Stores: Project DB (SQLite)  |  Execution Logs  |  Encrypted Env Files",
    sz=11, color=SLATE, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Implementation: System Setup
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "08", "Implementation \u2014 System Setup")
corner_dots(s)
page_number(s, 11)

card(s, Inches(0.8), Inches(1.4), Inches(5.6), Inches(3.0),
     "Application Initialization", [
         "FastAPI lifespan for startup/shutdown",
         "Stale devops-* container cleanup on start",
         "Asyncio task for 60s watchdog loop",
         "CORS + rate-limiting middleware",
         "SQLite table init via init_db()",
     ], ac=RUST, isz=14)

card(s, Inches(6.8), Inches(1.4), Inches(5.6), Inches(3.0),
     "Safety Configuration", [
         "DRY_RUN = true (log-only by default)",
         "ENABLE_LOCAL_EXECUTION = false",
         "ALLOWED_REPO_ROOTS = empty (demo)",
         "Pydantic-settings reads from .env",
         "Enforced at execution, not config time",
     ], ac=DUSTY_ROSE, isz=14)

card(s, Inches(0.8), Inches(4.7), Inches(11.6), Inches(2.3),
     "Docker Compose \u2014 One-Command Stack", [
         "Redis: redis:7-alpine with health check  |  Backend: FastAPI + Docker socket (port 3001)",
         "Worker: RQ background worker, shared backend-data volume  |  Frontend: Next.js (port 3000)",
         "Command: docker-compose up --build  \u2014  starts complete stack with health dependencies",
     ], ac=OLIVE, isz=15)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Command Interpretation
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "08", "Implementation \u2014 Command Interpretation")
corner_dots(s)
page_number(s, 12)

card(s, Inches(0.8), Inches(1.4), Inches(11.6), Inches(1.8),
     "Dual-Strategy Pipeline", [
         "1. User submits NL text  \u2192  2. LLM parse via Gemini (system prompt + 6 few-shot examples)  \u2192  3. Validate confidence",
         "4. If confidence < 0.4 OR LLM fails  \u2192  deterministic regex fallback  \u2192  5. Tagged with interpretation_method",
     ], ac=BROWN, isz=15)

card(s, Inches(0.8), Inches(3.5), Inches(5.6), Inches(3.5),
     "LLM Output Schema", [
         "action: deploy | rollback | status | ...",
         "version: semantic version or 'latest'",
         "environments: [staging, production, ...]",
         "post_steps: [run tests, smoke tests, ...]",
         "confidence: 0.0 \u2013 1.0 (threshold: 0.4)",
         "ai_reasoning: explanation of interpretation",
     ], ac=RUST, isz=14)

card(s, Inches(6.8), Inches(3.5), Inches(5.6), Inches(3.5),
     "Deterministic Regex Fallback", [
         "Activated on LLM failure or low confidence",
         "Regex patterns extract action, version, envs",
         "Result tagged: interpretation_method = 'deterministic'",
         "Frontend displays appropriate label",
         "System never fails on command parsing",
     ], ac=OLIVE, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Risk Assessment
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "08", "Implementation \u2014 Risk Assessment")
corner_dots(s)
page_number(s, 13)

card(s, Inches(0.8), Inches(1.4), Inches(5.6), Inches(2.5),
     "LLM Risk Advisor", [
         "Receives JSON deployment plan summary",
         "System prompt: 'reason as senior DevOps engineer'",
         "Returns: risk_level, risk_score (0\u2013100)",
         "  + warnings[], recommendations[]",
     ], ac=RUST, isz=14)

card(s, Inches(6.8), Inches(1.4), Inches(5.6), Inches(2.5),
     "Deterministic Fallback Scoring", [
         "Production environment: +30 points",
         "Missing tests: +20 | No smoke tests: +10",
         "Unknown action: +25 | No version: +10",
         "Total clamped to 100, mapped to level",
     ], ac=BROWN, isz=14)

# Risk level bar
add_rect(s, Inches(0.8), Inches(4.2), Inches(11.6), Inches(0.6), CREAM_DARK, border=BEIGE)
tf = tfb(s, Inches(1.0), Inches(4.25), Inches(11.2), Inches(0.5))
r = txt(tf, "Risk Levels:   ", sz=14, bold=True, color=BROWN, font=SERIF)
for label, threshold in [("LOW  (\u226430)", OLIVE), ("MEDIUM  (\u226460)", RGBColor(0xC4, 0x9A, 0x2F)),
                          ("HIGH  (\u226480)", RUST), ("CRITICAL  (\u2264100)", RGBColor(0xA8, 0x3C, 0x3C))]:
    run = r.add_run()
    run.text = f"   \u25CF {label}   "
    run.font.size = Pt(13)
    run.font.color.rgb = threshold
    run.font.bold = True
    run.font.name = SANS

card(s, Inches(0.8), Inches(5.1), Inches(11.6), Inches(1.9),
     "Plan Preview (shown before execution)", [
         "AI Confidence Score  \u2022  AI Reasoning  \u2022  Risk Level Badge  \u2022  Risk Score (0\u2013100)",
         "Warnings list  \u2022  Recommendations  \u2022  Explicit Approve / Cancel buttons",
         "No deployment executes without user clicking 'Approve' \u2014 mandatory safety gate",
     ], ac=OLIVE, isz=15)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — Orchestrator Pipeline
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "08", "Implementation \u2014 Six-Stage Pipeline")
corner_dots(s)
page_number(s, 14)

stages = [
    ("1", "Validate", "Path in whitelist\n.env present\nResolve workspace", OLIVE),
    ("2", "Git Pull", "GitHub projects\nSkipped for local\nBranch checkout", SLATE),
    ("3", "Build", "Docker Python SDK\nAuto-gen Dockerfile\nif none exists", RUST),
    ("4", "Stop Old", "Find latest deploy\nGraceful remove\nClear port conflicts", BROWN),
    ("5", "Run", "Launch container\nEnv file + ports\nEXPOSE > .env > default", DUSTY_ROSE),
    ("6", "Health Check", "5 retries, 2s interval\nRegister deployment\nStore LKG tag", OLIVE),
]

for i, (num, title, desc, col) in enumerate(stages):
    left = Inches(0.45) + i * Inches(2.1)
    add_rect(s, left, Inches(1.5), Inches(1.95), Inches(3.0), CREAM_DARK, border=BEIGE)
    add_sharp_rect(s, left, Inches(1.5), Inches(1.95), Pt(4), col)
    # Number circle
    c = circle(s, left + Inches(0.7), Inches(1.65), Inches(0.5), col)
    ctf = c.text_frame; cp = ctf.paragraphs[0]; cp.text = num
    cp.font.size = Pt(18); cp.font.bold = True; cp.font.color.rgb = WHITE
    cp.alignment = PP_ALIGN.CENTER; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Title
    tf = tfb(s, left + Inches(0.05), Inches(2.25), Inches(1.85), Inches(0.35))
    txt(tf, title, sz=14, bold=True, color=col, font=SERIF, align=PP_ALIGN.CENTER)
    # Desc
    tf2 = tfb(s, left + Inches(0.1), Inches(2.65), Inches(1.75), Inches(1.6))
    tf2.word_wrap = True
    txt(tf2, desc, sz=11, color=DARK_BROWN, align=PP_ALIGN.CENTER)

card(s, Inches(0.8), Inches(4.8), Inches(5.6), Inches(2.2),
     "Port Resolution", [
         "1. Dockerfile EXPOSE directive",
         "2. PORT= / APP_PORT= in .env",
         "3. Stack default: 8080 (Node), 8000 (Python)",
     ], ac=RUST, isz=14)

card(s, Inches(6.8), Inches(4.8), Inches(5.6), Inches(2.2),
     "Safety Controls", [
         "DRY_RUN: log-only, no Docker commands",
         "ENABLE_LOCAL_EXECUTION gate verified first",
         "Docker Python SDK only \u2014 no shell commands",
     ], ac=DUSTY_ROSE, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — Watchdog & Rollback
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "08", "Implementation \u2014 Watchdog & Rollback")
corner_dots(s)
page_number(s, 15)

card(s, Inches(0.8), Inches(1.4), Inches(5.6), Inches(3.2),
     "Watchdog Service", [
         "Runs every 60s via asyncio in FastAPI lifespan",
         "Iterates all projects, checks 'running' deployments",
         "Calls container_health_check() via Docker SDK",
         "Unhealthy: marks deployment as 'failed'",
         "Calls _rollback_to_last_known_good()",
     ], ac=DUSTY_ROSE, isz=14)

card(s, Inches(6.8), Inches(1.4), Inches(5.6), Inches(3.2),
     "Rollback Mechanism", [
         "Retrieves last_known_good_tag from project",
         "Decrypts and mounts environment file",
         "Starts new container with LKG image",
         "Creates new Deployment (status: running)",
         "Only rolls back if LKG differs from failed",
     ], ac=OLIVE, isz=14)

card(s, Inches(0.8), Inches(4.9), Inches(5.6), Inches(2.1),
     "Execution State Machine", [
         "queued \u2192 running | failed",
         "running \u2192 succeeded | failed | rolled_back",
         "succeeded / failed \u2192 rolled_back",
     ], ac=BROWN, isz=14)

card(s, Inches(6.8), Inches(4.9), Inches(5.6), Inches(2.1),
     "Plan State Machine", [
         "pending_approval \u2192 approved | cancelled",
         "approved \u2192 running | cancelled",
         "running \u2192 succeeded | failed (terminal)",
     ], ac=RUST, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — Chat Engine & LLM Abstraction
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "08", "Implementation \u2014 Chat & LLM Abstraction")
corner_dots(s)
page_number(s, 16)

card(s, Inches(0.8), Inches(1.4), Inches(5.6), Inches(2.8),
     "Conversation Engine", [
         "ChatSession + ChatMessage DB records",
         "Full history loaded on each turn",
         "Real-time project state in system prompt",
         "Parses <<<DEPLOY_INTENT>>> markers",
         "Auto-creates Plan on intent detection",
     ], ac=RUST, isz=14)

card(s, Inches(6.8), Inches(1.4), Inches(5.6), Inches(2.8),
     "Chat Frontend", [
         "Real-time messages with typing indicators",
         "Deployment intent cards in chat flow",
         "Inline plan preview + approval",
         "Session create / resume / list",
         "All API calls in lib/api.ts",
     ], ac=OLIVE, isz=14)

card(s, Inches(0.8), Inches(4.5), Inches(11.6), Inches(2.5),
     "LLM Provider Abstraction Layer", [
         "BaseLLMClient: complete() + complete_json() with 3-attempt JSON extraction strategy",
         "(1) Direct JSON parse  \u2192  (2) Markdown fence extraction  \u2192  (3) First {..} block  \u2192  (4) Retry with strict prompt",
         "GeminiClient: google-generativeai SDK  |  OpenAIClient: raw httpx (OpenAI + Ollama compatible)",
         "Factory selects provider via LLM_PROVIDER environment variable (GEMINI | OPENAI | OLLAMA)",
     ], ac=BROWN, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 17 — Screenshots: Dashboard & Chat
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "08", "Implementation \u2014 Screenshots (1/3)")
corner_dots(s)
page_number(s, 17)

# Dashboard — large placeholder
add_rect(s, Inches(0.8), Inches(1.4), Inches(11.6), Inches(5.5), CREAM_DARK, border=BEIGE)
add_sharp_rect(s, Inches(0.8), Inches(1.4), Inches(11.6), Pt(4), RUST)
tf = tfb(s, Inches(3.5), Inches(3.5), Inches(6.3), Inches(1.0))
tf.word_wrap = True
txt(tf, "Dashboard \u2014 Project List & Registration Panel", sz=18, color=RUST, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
txt(tf, "( insert screenshot )", sz=13, color=SLATE, italic=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 18 — Screenshots: Chat Interface
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "08", "Implementation \u2014 Screenshots (2/3)")
corner_dots(s)
page_number(s, 18)

# Chat — large placeholder
add_rect(s, Inches(0.8), Inches(1.4), Inches(11.6), Inches(5.5), CREAM_DARK, border=BEIGE)
add_sharp_rect(s, Inches(0.8), Inches(1.4), Inches(11.6), Pt(4), OLIVE)
tf = tfb(s, Inches(3.5), Inches(3.5), Inches(6.3), Inches(1.0))
tf.word_wrap = True
txt(tf, "Chat Interface \u2014 Multi-turn DevOps Conversation", sz=18, color=OLIVE, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
txt(tf, "( insert screenshot )", sz=13, color=SLATE, italic=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 19 — Screenshots: Plan Preview & Execution Logs
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "08", "Implementation \u2014 Screenshots (3/3)")
corner_dots(s)
page_number(s, 19)

# Two side-by-side, taller
for title, l, col in [
    ("Plan Preview \u2014 Risk Score & Approval Gate", Inches(0.8), BROWN),
    ("Execution Logs \u2014 6-Stage Pipeline Output", Inches(6.8), DUSTY_ROSE),
]:
    add_rect(s, l, Inches(1.4), Inches(5.6), Inches(5.5), CREAM_DARK, border=BEIGE)
    add_sharp_rect(s, l, Inches(1.4), Inches(5.6), Pt(4), col)
    tf = tfb(s, l + Inches(0.4), Inches(3.5), Inches(4.8), Inches(1.0))
    tf.word_wrap = True
    txt(tf, title, sz=15, color=col, bold=True, font=SERIF, align=PP_ALIGN.CENTER)
    txt(tf, "( insert screenshot )", sz=12, color=SLATE, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 20 — Testing: Unit & Integration
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "08", "Testing \u2014 Unit & Integration")
corner_dots(s)
page_number(s, 20)

card(s, Inches(0.8), Inches(1.4), Inches(5.6), Inches(2.8),
     "Unit Tests (pytest + anyio)", [
         "test_demo_analysis.py: stack detection + Dockerfile gen",
         "Verifies Node.js \u2192 'node' stack, multi-stage build",
         "test_m3_orchestrator.py: Docker pipeline test",
         "Mocks DB/env, real Docker operations",
         "Auto-skips when Docker unavailable",
     ], ac=RUST, isz=14)

card(s, Inches(6.8), Inches(1.4), Inches(5.6), Inches(2.8),
     "Integration Tests", [
         "Dual-strategy fallback verification",
         "LLM mocked to simulate failures",
         "Verifies interpretation_method = 'deterministic'",
         "Low confidence \u2192 threshold logic validation",
         "Risk advisor LLM/deterministic interaction",
     ], ac=OLIVE, isz=14)

card(s, Inches(0.8), Inches(4.5), Inches(11.6), Inches(2.5),
     "Functional (E2E) Scenarios", [
         "Happy path: register \u2192 upload .env \u2192 NL command \u2192 review plan \u2192 approve \u2192 verify container",
         "Risk escalation: 'deploy to production no tests' \u2192 verify risk > 60 and level = high/critical",
         "Fallback: invalid LLM_PROVIDER \u2192 deterministic parser handles gracefully",
         "Watchdog: manually stop container \u2192 wait 60s \u2192 verify auto-rollback to LKG image",
         "Chat: 'deploy latest to staging' in chat \u2192 Plan auto-created and linked to session",
     ], ac=BROWN, isz=14)


# ════════════════════════════════════════════════════════════════
# SLIDE 21 — Testing: Performance & Security
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "08", "Testing \u2014 Performance & Security")
corner_dots(s)
page_number(s, 21)

card(s, Inches(0.8), Inches(1.4), Inches(5.6), Inches(2.5),
     "Performance", [
         "LLM command parse: 3\u20136 seconds (Gemini)",
         "Deterministic fallback: < 50 ms",
         "RQ job enqueue: < 100 ms",
         "Cached Docker build: < 10 seconds",
     ], ac=OLIVE, isz=15)

card(s, Inches(6.8), Inches(1.4), Inches(5.6), Inches(2.5),
     "Security", [
         "Dry-run: no containers with DRY_RUN=true",
         "Execution gate: clean fail when EXEC=false",
         "Path traversal: rejected outside whitelist",
         "Rate limiting: HTTP 429 after 30 req/min/IP",
     ], ac=DUSTY_ROSE, isz=15)

card(s, Inches(0.8), Inches(4.2), Inches(11.6), Inches(2.8),
     "User Acceptance Testing", [
         "Scenario: developer with no Docker experience deploys a Node.js app to staging",
         "Types 'ship the app to staging and run tests' in chat \u2192 reviews AI plan \u2192 approves \u2192 observes logs",
         "Follow-up: 'what was the last deployment status?' \u2192 contextual response from real project state",
         "Feedback: NL interface significantly lowered perceived complexity of deployment",
     ], ac=BROWN, isz=15)


# ════════════════════════════════════════════════════════════════
# SLIDE 20 — Conclusion
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "09", "Conclusion")
corner_dots(s)
page_number(s, 22)

card(s, Inches(0.8), Inches(1.4), Inches(11.6), Inches(2.0),
     "What We Achieved", [
         "Production-grade deployment automation made accessible and safe through LLM agents + deterministic fallbacks",
         "Successfully bridges the gap between natural language intent and complex infrastructure operations",
         "Dual-strategy ensures reliability: LLM primary, deterministic fallback \u2014 system never fails to respond",
     ], ac=OLIVE, isz=15)

achievements = [
    ("Intelligent\nUnderstanding", "LLM interprets intent;\nreasons, not just executes", RUST),
    ("Risk\nAwareness", "Deployment intelligence\nabsent from CI/CD tools", OLIVE),
    ("Mandatory\nSafety Gates", "Approval + dry-run prevent\naccidental destruction", BROWN),
    ("Self-Healing", "Watchdog auto-rollback\nfor container failures", DUSTY_ROSE),
    ("Conversational\nAssistant", "Multi-turn chat with\nreal-time project state", SLATE),
]

for i, (title, desc, col) in enumerate(achievements):
    left = Inches(0.5) + i * Inches(2.55)
    add_rect(s, left, Inches(3.8), Inches(2.35), Inches(3.0), CREAM_DARK, border=BEIGE)
    add_sharp_rect(s, left, Inches(3.8), Inches(2.35), Pt(4), col)
    tf = tfb(s, left + Inches(0.1), Inches(3.95), Inches(2.15), Inches(0.75))
    tf.word_wrap = True
    txt(tf, title, sz=15, bold=True, color=col, font=SERIF, align=PP_ALIGN.CENTER)
    tf2 = tfb(s, left + Inches(0.1), Inches(4.8), Inches(2.15), Inches(1.6))
    tf2.word_wrap = True
    txt(tf2, desc, sz=12, color=DARK_BROWN, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 21 — Future Enhancements
# ════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "10", "Future Enhancements")
corner_dots(s)
page_number(s, 23)

phases = [
    ("Cloud Provider\nIntegration", RUST, [
        "Native Vercel & Render APIs",
        "Unified deploy status endpoint",
        "Cross-platform orchestration",
        "Normalized status schema",
    ]),
    ("Production\nHardening", OLIVE, [
        "PostgreSQL + Alembic migrations",
        "WebSocket log streaming",
        "Horizontal RQ worker scaling",
        "Git branch-aware deployments",
    ]),
    ("Observability &\nAnalytics", BROWN, [
        "Deployment history timeline",
        "Prometheus + Grafana export",
        "Slack / email rollback alerts",
        "Success rate & MTTR metrics",
    ]),
    ("Advanced AI\nCapabilities", DUSTY_ROSE, [
        "Multi-project dependency reasoning",
        "SOP ingestion via Chroma vector DB",
        "Predictive risk from history",
        "True RAG with deployment SOPs",
    ]),
]

for i, (title, color, items) in enumerate(phases):
    left = Inches(0.6) + i * Inches(3.15)
    add_rect(s, left, Inches(1.5), Inches(2.95), Inches(5.3), CREAM_DARK, border=BEIGE)
    add_sharp_rect(s, left, Inches(1.5), Inches(2.95), Pt(4), color)
    tf = tfb(s, left + Inches(0.15), Inches(1.7), Inches(2.65), Inches(0.8))
    tf.word_wrap = True
    txt(tf, title, sz=16, bold=True, color=color, font=SERIF, align=PP_ALIGN.CENTER)
    thin_line(s, left + Inches(0.4), Inches(2.55), Inches(2.15), color=BEIGE)
    bullet_list(s, items, left + Inches(0.15), Inches(2.75), Inches(2.65), Inches(3.8),
                sz=13, col=DARK_BROWN, bcol=TAN, sp=Pt(8))


# ════════════════════════════════════════════════════════════════
# SLIDE 22 — Thank You
# ════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s)

# Decorative circles
circle(s, Inches(9.8), Inches(-1), Inches(4.5), CREAM_DARK, border=BEIGE)
circle(s, Inches(-0.5), Inches(5.5), Inches(3), CREAM_DARK, border=BEIGE)
circle(s, Inches(11), Inches(5.8), Inches(2), CREAM_DARK, border=BEIGE)

# Vertical accent line
add_sharp_rect(s, Inches(1.5), Inches(1.8), Pt(1.5), Inches(4), TAN)

# Content
tf = tfb(s, Inches(2.2), Inches(2.0), Inches(8), Inches(1.0))
txt(tf, "Thank You", sz=52, color=CHARCOAL, font=SERIF)

thin_line(s, Inches(2.2), Inches(3.1), Inches(3.5), color=TAN, thick=Pt(1))

tf2 = tfb(s, Inches(2.2), Inches(3.4), Inches(8), Inches(2.5))
tf2.word_wrap = True
txt(tf2, "AI DevOps Commander", sz=22, bold=True, color=BROWN, font=SERIF, after=Pt(4))
txt(tf2, "Natural Language Driven Deployment Automation with Agentic AI", sz=14, color=SLATE, italic=True, after=Pt(20))
txt(tf2, "Prince Parmar Singh (24100443)", sz=16, color=CHARCOAL, font=SERIF, after=Pt(4))
txt(tf2, "Guide: Dr. Deepika MP, Assistant Professor", sz=14, color=SLATE, after=Pt(4))
txt(tf2, "Department of Computer Applications, CUSAT", sz=13, color=SLATE)

# Question prompt
tf3 = tfb(s, Inches(2.2), Inches(5.8), Inches(8), Inches(0.5))
txt(tf3, "Questions & Discussion", sz=18, color=BROWN, bold=False, font=SERIF, italic=True)


# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
out_path = "AI_DevOps_Commander_Presentation_v2.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
